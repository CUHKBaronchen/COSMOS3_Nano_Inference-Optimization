#!/usr/bin/env python3
from pathlib import Path
from xml.sax.saxutils import escape
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import BaseDocTemplate, Frame, PageTemplate, Paragraph, Spacer, Table, TableStyle, Image as RImage, PageBreak

ROOT=Path('/root/autodl-tmp')
OUT=ROOT/'cosmos3_final_delivery'
GEN=OUT/'assets/generated'
FONT=OUT/'assets/fonts/NotoSansCJKsc-Regular.otf'
PDF_REG=OUT/'assets/fonts/NotoSansSC-Regular.ttf'
PDF_BOLD=OUT/'assets/fonts/NotoSansSC-Bold.ttf'
GEN.mkdir(parents=True,exist_ok=True)
C={'navy':'#0B1426','navy2':'#13213A','blue':'#36A3FF','green':'#4FD1A1','orange':'#FFB454','red':'#FF6B6B','white':'#F7FAFC','muted':'#AAB8CF','ink':'#182230','grid':'#2A3A55'}
B0=481.361; P2=418.109; P3=359.067; P4=344.61192417144775
paths={
'Baseline video':ROOT/'cosmos3_t2v_baseline/outputs/p0_baseline_20260613T115211Z/t2v/vision.mp4',
'Baseline benchmark':ROOT/'cosmos3_t2v_baseline/outputs/p0_baseline_20260613T115211Z/benchmark.json',
'Exact initial noise':ROOT/'cosmos3_t2v_baseline/outputs/p0_capture_20260613T114518Z/t2v/initial_noise.safetensors',
'Profiler trace':ROOT/'cosmos3_t2v_baseline/outputs/p0_capture_20260613T114518Z/t2v/profile.json.gz',
'P2 output':ROOT/'cosmos3_t2v_p2/outputs/cfg_interval_800_1000_20260614T130154Z/t2v/vision.mp4',
'P3 output':ROOT/'cosmos3_t2v_p3/outputs/fp8_steps_35_20260614T160455Z/t2v/vision.mp4',
'P4 canonical output':ROOT/'cosmos3_t2v_p4/outputs/residual_steps_35_reuse_11_13_15_17_19_21_23_balanced_20260615T0254Z/t2v/vision.mp4',
'P4 repeat benchmark':ROOT/'cosmos3_t2v_p4/outputs/residual_steps_35_reuse_11_13_15_17_19_21_23_balanced_steady_repeat3_20260615T1418Z/benchmark.json',
'Quality manifest':ROOT/'cosmos3_t2v_quality/QUALITY_MANIFEST.json',
'P4 status':ROOT/'cosmos3_t2v_p4/P4_STATUS.md'}
videos={
'B0 vs P2 CFG':ROOT/'cosmos3_t2v_p2/B0_35_left_vs_CFG800_1000_right_same_noise.mp4',
'B0 vs P3 FP8+CFG':ROOT/'cosmos3_t2v_p3/B0_BF16_left_vs_P3_FP8_CFG_right_same_noise.mp4',
'P3 vs P4 Balanced':ROOT/'cosmos3_t2v_p4/P3_left_vs_P4_balanced_right_same_noise.mp4',
'human_handoff P3 vs P4':ROOT/'cosmos3_t2v_quality/human_handoff_P3_left_vs_P4_right_same_noise.mp4',
'rally_tracking P3 vs P4':ROOT/'cosmos3_t2v_quality/rally_tracking_P3_left_vs_P4_right_same_noise.mp4',
'tabletop_relations P3 vs P4':ROOT/'cosmos3_t2v_quality/tabletop_relations_P3_left_vs_P4_right_same_noise.mp4',
'Aggressive boundary':ROOT/'cosmos3_t2v_quality/human_handoff_P4_balanced_left_vs_aggressive_right_same_noise.mp4'}
contacts={
'cfg':ROOT/'cosmos3_t2v_p2/B0_35_left_vs_CFG800_1000_right_contact_sheet.png',
'fp8':ROOT/'cosmos3_t2v_p3/B0_BF16_left_vs_P3_FP8_CFG_right_contact_sheet.png',
'p4':ROOT/'cosmos3_t2v_p4/P3_left_vs_P4_balanced_right_contact_sheet.png',
'human':ROOT/'cosmos3_t2v_quality/human_handoff_P3_left_vs_P4_right_contact_sheet.png',
'rally':ROOT/'cosmos3_t2v_quality/rally_tracking_P3_left_vs_P4_right_contact_sheet.png',
'table':ROOT/'cosmos3_t2v_quality/tabletop_relations_P3_left_vs_P4_right_contact_sheet.png',
'aggr':ROOT/'cosmos3_t2v_quality/human_handoff_P4_balanced_left_vs_aggressive_right_contact_sheet.png'}

def F(n): return ImageFont.truetype(str(FONT),n)
def wrap(draw,text,font,maxw):
    out=[]
    for para in text.split('\n'):
        cur=''
        for ch in para:
            if draw.textbbox((0,0),cur+ch,font=font)[2]<=maxw or not cur: cur+=ch
            else: out.append(cur);cur=ch
        out.append(cur)
    return out
def draw_lines(draw,text,xy,font,fill,maxw,space=8):
    x,y=xy
    for line in wrap(draw,text,font,maxw): draw.text((x,y),line,font=font,fill=fill);y+=font.size+space
    return y
def fit_image(src,dst,maxw=2200):
    im=Image.open(src).convert('RGB')
    if im.width>maxw: im=im.resize((maxw,round(im.height*maxw/im.width)),Image.Resampling.LANCZOS)
    im.save(dst,quality=94)
def chart_e2e():
    im=Image.new('RGB',(1400,760),'white');d=ImageDraw.Draw(im)
    d.text((60,35),'端到端时延：逐项优化消融',font=F(44),fill=C['ink'])
    vals=[('B0 BF16',B0,C['red']),('P2 CFG',P2,C['orange']),('P3 + FP8',P3,C['blue']),('P4 + Cache',P4,C['green'])]
    x0,y0,pw,ph=100,130,1200,500
    for t in range(0,501,100):
        y=y0+ph-t/500*ph;d.line((x0,y,x0+pw,y),fill='#D8E1EA',width=2);d.text((25,y-13),str(t),font=F(22),fill='#667085')
    for i,(name,val,col) in enumerate(vals):
        x=180+i*285;h=val/500*ph;y=y0+ph-h
        d.rounded_rectangle((x,y,x+185,y0+ph),radius=16,fill=col)
        d.text((x+92,y-34),f'{val:.3f}s',font=F(27),fill=C['ink'],anchor='mm');d.text((x+92,675),name,font=F(23),fill=C['ink'],anchor='mm')
    d.text((60,720),f'最终加速 {B0/P4:.4f}x；端到端时延下降 {(1-P4/B0)*100:.2f}%',font=F(27),fill='#155EEF')
    im.save(GEN/'e2e.png')
def chart_breakdown():
    im=Image.new('RGB',(1500,820),'white');d=ImageDraw.Draw(im);d.text((60,35),'各环节时间 Breakdown（代表性 warmed run）',font=F(42),fill=C['ink'])
    rows=[('B0',481.494,[231.918,220.726,18.071,5.963]),('P2',418.109,[231.900,157.650,18.073,5.962]),('P3',359.067,[197.415,133.020,18.070,5.944]),('P4',344.400,[182.772,133.025,18.065,5.943])]
    names=['Conditional','Unconditional','VAE decode','Prepare','Other'];cols=[C['blue'],C['orange'],C['green'],'#36D1DC','#9A7BFF'];scale=2.3
    for i,(name,total,parts) in enumerate(rows):
        y=155+i*120;d.text((65,y+22),name,font=F(28),fill=C['ink']);x=220
        for val,col in zip(parts+[total-sum(parts)],cols):
            w=max(2,val*scale);d.rectangle((x,y,x+w,y+65),fill=col)
            if w>105:d.text((x+w/2,y+32),f'{val:.1f}',font=F(20),fill='white',anchor='mm')
            x+=w
        d.text((x+12,y+32),f'{total:.1f}s',font=F(24),fill=C['ink'],anchor='lm')
    for i,(n,col) in enumerate(zip(names,cols)):
        x=190+(i%3)*400;y=660+(i//3)*55;d.rectangle((x,y,x+34,y+28),fill=col);d.text((x+48,y+14),n,font=F(21),fill=C['ink'],anchor='lm')
    d.text((60,775),'P2 减少 unconditional；P3/P4 继续压缩双分支主干。',font=F(26),fill='#155EEF');im.save(GEN/'breakdown.png')
def chart_cache():
    vals=[(7,.069336),(9,.059814),(11,.054443),(13,.050537),(15,.050293),(17,.050537),(19,.052734),(21,.053711),(23,.055176),(27,.064941),(29,.079590)]
    im=Image.new('RGB',(1400,760),'white');d=ImageDraw.Draw(im);d.text((60,35),'P4 校准：跨 timestep 输入变化呈 U 形',font=F(42),fill=C['ink'])
    x0,y0,pw,ph=120,135,1160,470;lo,hi=.045,.085;pts=[]
    for i in range(5):
        v=lo+i*(hi-lo)/4;y=y0+ph-(v-lo)/(hi-lo)*ph;d.line((x0,y,x0+pw,y),fill='#D8E1EA',width=2);d.text((25,y-12),f'{v:.3f}',font=F(20),fill='#667085')
    for step,v in vals:pts.append((x0+(step-7)/22*pw,y0+ph-(v-lo)/(hi-lo)*ph))
    d.line(pts,fill=C['blue'],width=6)
    for (step,v),(x,y) in zip(vals,pts):d.ellipse((x-9,y-9,x+9,y+9),fill=C['green'] if 11<=step<=23 else C['orange']);d.text((x,y-28),str(step),font=F(20),fill=C['ink'],anchor='mm')
    th=.0545;y=y0+ph-(th-lo)/(hi-lo)*ph;d.line((x0,y,x0+pw,y),fill=C['red'],width=3);d.text((1270,y-18),'threshold 0.0545',font=F(21),fill=C['red'],anchor='ra')
    d.text((60,690),'首尾 5 步与 CFG 切换点受保护；候选非相邻；阈值失败自动完整刷新。',font=F(25),fill='#155EEF');im.save(GEN/'cache.png')
def chart_stability():
    vals=[344.400141,344.662216,344.611924];im=Image.new('RGB',(1400,700),'white');d=ImageDraw.Draw(im);d.text((60,35),'P4 三进程稳态复测',font=F(42),fill=C['ink'])
    x0,y0,pw,ph=150,140,1100,360;lo,hi=344.2,344.85;pts=[]
    for i in range(4):v=lo+i*(hi-lo)/3;y=y0+ph-(v-lo)/(hi-lo)*ph;d.line((x0,y,x0+pw,y),fill='#D8E1EA',width=2);d.text((35,y-12),f'{v:.2f}',font=F(20),fill='#667085')
    for i,v in enumerate(vals):pts.append((x0+100+i*450,y0+ph-(v-lo)/(hi-lo)*ph))
    d.line(pts,fill=C['blue'],width=6)
    for i,((x,y),v) in enumerate(zip(pts,vals),1):d.ellipse((x-13,y-13,x+13,y+13),fill=C['green']);d.text((x,y-40),f'{v:.3f}s',font=F(25),fill=C['ink'],anchor='mm');d.text((x,550),f'Run {i}',font=F(24),fill=C['ink'],anchor='mm')
    d.text((60,635),'Median 344.612s | Range 0.262s (0.076%) | CV≈0.033% | VRAM 41,233 MiB',font=F(25),fill='#155EEF');im.save(GEN/'stability.png')
def chart_quality():
    im=Image.new('RGB',(1400,720),'white');d=ImageDraw.Draw(im);d.text((60,35),'人工质量门：同初始噪声、完整视频定性审核',font=F(40),fill=C['ink'])
    headers=['Prompt / 风险','B0 vs P3','P3 vs P4','Balanced vs Aggressive'];rows=[['Anchor 工业铸造','PASS','PASS','背景细节下降'],['human_handoff','PASS','PASS / P4略优','FAIL / 边缘模糊'],['rally_tracking','P3略优','PASS / P4略优','对比减弱，物理感更真'],['tabletop_relations','PASS','PASS / P4略优','PASS / 物理更合理']]
    xs=[60,430,690,980,1340];y0=130;rh=110
    for i,h in enumerate(headers):d.rectangle((xs[i],y0,xs[i+1],y0+65),fill='#17365D');d.text(((xs[i]+xs[i+1])/2,y0+32),h,font=F(21),fill='white',anchor='mm')
    for r,row in enumerate(rows):
        y=y0+65+r*rh
        for i,t in enumerate(row):
            d.rectangle((xs[i],y,xs[i+1],y+rh),fill='#F3F6FA' if r%2==0 else '#E8EEF5',outline='#C8D2E0');draw_lines(d,t,(xs[i]+12,y+16),F(20),C['red'] if 'FAIL' in t else C['ink'],xs[i+1]-xs[i]-24,5)
    d.text((60,660),'Balanced 通过 Anchor + 3/3；Aggressive 因人物边缘模糊淘汰。',font=F(25),fill='#155EEF');im.save(GEN/'quality.png')

def assets():
    chart_e2e();chart_breakdown();chart_cache();chart_stability();chart_quality()
    for k,v in contacts.items():fit_image(v,GEN/f'{k}.jpg')

def write_sources():
    L=['# Cosmos3-Nano 16B 单卡 720p 5s T2V 免训练推理加速报告','', '**最终结果：344.612 s；相对本机同协议 B0 为 1.3968×；端到端时延下降 28.41%；Anchor + 3/3 高风险 Prompt 通过。**','', '## 摘要','', '本项目在 NVIDIA Cosmos Framework 上，对 Cosmos3-Nano 16B 的文本生成视频路径进行单卡、720p、121 帧、24 FPS、35-step UniPC 推理优化。全流程不训练、不蒸馏、不更改 VAE。最终方案由后期 CFG 分支裁剪、generation pathway 选择性 FP8、动态阈值 conditional residual cache 组成。','', '## 固定协议','', '- GPU：1× NVIDIA H20 96GB','- 输出：1280×720，121 frames，24 FPS，首尾帧跨度 5.000s','- Sampler：UniPC 35 steps，shift 10，CFG 6','- Noise SHA-256：`44122c662d3bd10659534486aa1d455a22ffe6d0c24ffb0dbb5fe7f0553d4b61`','- 质量：同初始噪声完整视频人工审核，不以 PSNR/SSIM 替代','', '## 优化结论','', '1. P1：Static compile、CUDA Graph、text K/V cache 均无可计入收益，保留为负结果。','2. P2：最后 10 个低噪声 step 跳过 unconditional CFG，NFE 35/35→35/25，E2E 418.109s。','3. P3：选择性量化第 4–31 层 140 个 generation Linear，E2E median 359.067s。','4. P4：conditional 层 8–27 动态 residual reuse，候选 11–23，阈值 0.0545，E2E median 344.612s。','', '## 消融结果','', '|配置|Steps|Cond/Uncond|Precision|Cache|E2E|vs B0|质量|','|---|---:|---:|---|---|---:|---:|---|','|B0|35|35/35|BF16|Off|481.361s median|1.0000×|Reference|','|P2|35|35/25|BF16|Off|418.109s|1.1513×|Anchor PASS|','|P3|35|35/25|Selective FP8|Off|359.067s median|1.3406×|Anchor+3/3 PASS|','|P4 Balanced|35|35/25|Selective FP8|Dynamic residual|344.612s median|1.3968×|Anchor+3/3 PASS|','|P4 Aggressive|35|35/25|Selective FP8|Aggressive|335.881s|1.4331×|Reject: 边缘模糊|','', '## 官方参考口径','', 'NVIDIA 当前 H20 PyTorch 720p/1 为 931.39s，但官方说明标准视频 workload 通常为 189 frames；本项目为 121 frames，且计时范围不同。因此官方值只作背景，不用于计算本项目加速比。','', '官方链接：https://github.com/NVIDIA/cosmos/blob/main/inference_benchmarks.md#text-to-video-t2v','', '## 绝对路径证据索引','']
    for k,v in paths.items():L.append(f'- **{k}**: `{v}`')
    L+=['','## 同噪声完整视频','']
    for k,v in videos.items():L.append(f'- **{k}**: `{v}`')
    L+=['','## Commit','', '- Framework base: `a5ae92b7d7aab100e2f5e96c44788adfce26331c`','- NVIDIA cosmos reference: `1fe7e3be1687d797392b0e82ff6fe6296638b49f`','', '## 结论','', '最终 P4 Balanced 在保持 35-step UniPC 和严格同噪声质量门的前提下实现 1.3968× E2E 加速。报告公开负结果和 Aggressive 失败案例，保证结果可审计、可复现。']
    (OUT/'COSMOS3_T2V_FINAL_REPORT.md').write_text('\n'.join(L)+'\n',encoding='utf-8')
    I=['# Evidence Index','', '所有路径均为服务器绝对路径。','']+[f'- `{k}`: `{v}`' for k,v in paths.items()]+['','## Videos','']+[f'- `{k}`: `{v}`' for k,v in videos.items()]
    (OUT/'EVIDENCE_INDEX.md').write_text('\n'.join(I)+'\n',encoding='utf-8')

class Doc(BaseDocTemplate):
    def __init__(self,name):
        super().__init__(name,pagesize=A4,leftMargin=18*mm,rightMargin=18*mm,topMargin=18*mm,bottomMargin=16*mm,title='Cosmos3 T2V Final Report')
        self.addPageTemplates(PageTemplate(id='p',frames=Frame(18*mm,17*mm,A4[0]-36*mm,A4[1]-34*mm),onPage=self.hf))
    def hf(self,c,d):
        c.saveState();c.setFont('Noto',8);c.setFillColor(colors.HexColor('#667085'));c.drawString(18*mm,A4[1]-12*mm,'Cosmos3-Nano 16B | Single-H20 720p 5s T2V');c.drawRightString(A4[0]-18*mm,10*mm,str(d.page));c.restoreState()
def rimg(p,w=172):
    im=Image.open(p);return RImage(str(p),width=w*mm,height=w*im.height/im.width*mm)
def build_pdf():
    pdfmetrics.registerFont(TTFont('Noto',str(PDF_REG)));pdfmetrics.registerFont(TTFont('NotoB',str(PDF_BOLD)))
    body=ParagraphStyle('body',fontName='Noto',fontSize=9.2,leading=14,textColor=colors.HexColor(C['ink']),spaceAfter=5)
    h1=ParagraphStyle('h1',fontName='NotoB',fontSize=18,leading=24,textColor=colors.HexColor('#0B4AA2'),spaceBefore=7,spaceAfter=7)
    h2=ParagraphStyle('h2',fontName='NotoB',fontSize=13,leading=18,textColor=colors.HexColor('#155EEF'),spaceBefore=6,spaceAfter=5)
    title=ParagraphStyle('title',fontName='NotoB',fontSize=26,leading=36,alignment=TA_CENTER,textColor=colors.HexColor('#0B4AA2'))
    sub=ParagraphStyle('sub',fontName='Noto',fontSize=12,leading=19,alignment=TA_CENTER,textColor=colors.HexColor('#475467'))
    small=ParagraphStyle('small',fontName='Noto',fontSize=7.2,leading=10,textColor=colors.HexColor('#475467'))
    S=[]
    def P(t,s=body):return Paragraph(t,s)
    def H(t):S.append(P(t,h1))
    def H2(t):S.append(P(t,h2))
    def TB(rows,widths,fs=7.3):
        st=ParagraphStyle('cell',parent=body,fontSize=fs,leading=fs+3)
        data=[[P(escape(str(x)),st) for x in r] for r in rows];t=Table(data,colWidths=widths,repeatRows=1)
        t.setStyle(TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#17365D')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),.35,colors.HexColor('#B8C4D4')),('VALIGN',(0,0),(-1,-1),'MIDDLE'),('ROWBACKGROUNDS',(0,1),(-1,-1),[colors.white,colors.HexColor('#F6F8FB')]),('PADDING',(0,0),(-1,-1),4)]));S.extend([t,Spacer(1,3*mm)])
    def bullets(items):
        for x in items:S.append(P('• '+escape(x)))
    S.extend([Spacer(1,30*mm),P('Cosmos3-Nano 16B',title),P('单卡 720p 5s T2V 免训练推理加速报告',title),Spacer(1,8*mm),P('NVIDIA H20 · 35-step UniPC · Same-noise Human Quality Gate',sub),Spacer(1,12*mm)])
    TB([['最终 E2E','端到端加速','时延下降','质量门'],['344.612 s','1.3968×','28.41%','Anchor + 3/3 PASS']],[43*mm]*4,10)
    S.extend([Spacer(1,10*mm),P('最终组合：后期 CFG 裁剪 + generation-selective FP8 + dynamic conditional residual cache。保持默认 35 步，不训练、不蒸馏、不修改 VAE。',sub),PageBreak()])
    H('1. 执行摘要');S.append(P('本项目针对 Cosmos3-Nano 16B 文本生成视频路径，在单张 NVIDIA H20 上固定 1280×720、121 帧、24 FPS、UniPC 35 steps、shift 10、CFG 6。所有核心视觉比较使用完全相同的 initial latent noise，最终质量由完整视频人工定性审核。'));S.append(rimg(GEN/'e2e.png'))
    TB([['指标','B0','P4 Balanced','结论'],['E2E median','481.361s','344.612s','1.3968×'],['时延下降','—','136.749s','28.41%'],['Peak VRAM','45,871 MiB','41,233 MiB','-4,638 MiB'],['稳态波动','—','Range 0.262s','CV≈0.033%'],['质量','Reference','Anchor + 3/3 PASS','Aggressive reject']],[38*mm,38*mm,45*mm,48*mm])
    H('2. 任务映射与固定实验协议');TB([['要求','落实'],['单卡 720p 5s','1×H20；1280×720；121f；24 FPS'],['免训练','CFG interval、FP8、residual reuse'],['Diffusion 加速','保持 35 steps，减少后期 unconditional NFE'],['Kernel','确认 FA3；compile/graph 消融'],['量化','140 个 generation Linear FP8'],['质量','同噪声完整视频人工审核'],['数据','E2E、breakdown、GPU telemetry']],[45*mm,125*mm]);bullets(['Noise SHA-256: 44122c662d3bd10659534486aa1d455a22ffe6d0c24ffb0dbb5fe7f0553d4b61','每个正式进程先完整 warmup，再记录 measured run','Framework base commit: a5ae92b7d7aab100e2f5e96c44788adfce26331c'])
    H('3. Baseline 与瓶颈归因');S.append(P('B0 代表性 warmed run 为 481.494s，三次中位数 481.361s。Conditional / unconditional MoT forward 合计占 94.01% E2E；VAE decode 仅 3.75%。Profiler 记录 288 次 flash_attn_3::_flash_attn_forward，确认 FA3 已生效。'));S.append(rimg(GEN/'breakdown.png'))
    TB([['Stage','B0','P2','P3','P4 rep.'],['Prepare','5.963','5.962','5.944','5.943'],['Conditional','231.918','231.900','197.415','182.772'],['Unconditional','220.726','157.650','133.020','133.025'],['VAE decode','18.071','18.073','18.070','18.065'],['E2E','481.494','418.109','359.067','344.400']],[34*mm]*5)
    H('4. P1：负结果同样进入报告');TB([['候选','结果','决定'],['Static torch.compile','481.142s；1.00046×','噪声级，不合入'],['CUDA Graph smoke','慢0.23%；VRAM +3,068MiB','停止'],['Text K/V cache','2-step 1.0037×','收益不足'],['FA3','288 profiler calls','已在 baseline 生效']],[45*mm,72*mm,53*mm]);S.append(P('结论：当前 workload 由大型 MoT 算子计算主导，launch/Python 不是主要瓶颈。'))
    H('5. P2：35-step CFG interval');S.append(P('设置 guidance_interval=[800,1000]。Scheduler 和 conditional 都保持 35 次；最后 10 个低噪声 step 跳过 unconditional，NFE 35/35→35/25。E2E 418.109s，相对 B0 1.1513×。'));TB([['Metric','B0','P2','变化'],['Conditional','231.918','231.900','不变'],['Unconditional','220.726','157.650','-63.077s'],['Denoise','458.682','395.580','1.1595×'],['E2E','481.494','418.109','1.1516×']],[42*mm]*4);S.append(rimg(GEN/'cfg.jpg'));S.append(P('左 B0 / 右 P2，同噪声五时间点。30/28-step 仅作为敏感度消融，不进入最终主线。',small))
    H('6. P3：Generation-selective FP8');bullets(['第4–31层，共140个 generation Linear','q_proj_moe_gen、o_proj_moe_gen、MLP gate/up/down','K/V、text branch、首尾4层、Norm、RoPE、Attention、VAE保持BF16','三进程 E2E median 359.067s，CV 0.0127%']);TB([['Metric','P2','P3','收益'],['Conditional','231.900','197.415','1.1747×'],['Unconditional','157.650','133.020','1.1852×'],['Denoise','395.580','336.448','1.1758×'],['E2E','418.109','359.067','1.1644×']],[42*mm]*4);S.append(rimg(GEN/'fp8.jpg'))
    H('7. P4：Conditional residual cache');S.append(P('仅在 conditional 分支第8–27层工作。进入第8层前比较当前与上一刷新步 generation hidden state 的 relative-L1；低于0.0545才复用，否则自动完整刷新。候选11,13,15,17,19,21,23；首尾5步和CFG transition受保护；候选必须非相邻。'));S.append(rimg(GEN/'cache.png'));TB([['机制','设置'],['候选 steps','11,13,15,17,19,21,23'],['层范围','8–27'],['阈值','relative-L1≤0.0545'],['保护','first5 / last5 / step25'],['回退','超阈值完整执行']],[55*mm,115*mm]);S.append(P('每个正式进程命中5步，共跳过100个 conditional decoder block。Conditional 197.415→182.772s；unconditional 与 VAE 基本不变。'))
    H('8. 最终稳态性能');S.append(rimg(GEN/'stability.png'));TB([['Run','E2E','Cond','Uncond','VRAM','Accepted'],['Original','344.400','182.772','133.025','41,233','11,13,15,17,21'],['Repeat2','344.662','182.790','133.031','41,233','13,15,17,19,23'],['Repeat3','344.612','182.797','133.046','41,233','11,13,15,17,21']],[27*mm,25*mm,25*mm,25*mm,27*mm,42*mm]);S.append(P('FP8 数值波动会改变阈值边界的具体命中 step，但命中数、E2E 和显存稳定。两轮 WARMUP=0 冷启动 353.741/354.725s 不混入稳态中位数。'))
    H('9. 同噪声人工质量门');S.append(rimg(GEN/'quality.png'));S.append(P('审核 prompt 语义、主体数量与关系、动作、镜头运动、边缘/纹理、模糊、噪点、闪烁和物理合理性。联系图用于筛查，最终结论来自完整 MP4。'))
    for k,t in [('p4','Anchor：P3（左）vs P4 Balanced（右）'),('human','人物交接'),('rally','快速运动与跟拍'),('table','多物体空间关系')]:H2(t);S.append(rimg(GEN/f'{k}.jpg'))
    H('10. 为什么不选择更快的 Aggressive');S.append(P('Aggressive anchor E2E 335.881s，仅比 Balanced 快约2.6%，但 human_handoff 出现可感知人物/物体边缘模糊，因此严格淘汰。'));S.append(rimg(GEN/'aggr.jpg'));TB([['Prompt','观察','结论'],['human_handoff','边缘更软/模糊','FAIL'],['rally_tracking','线条/颜色对比减弱，物理感更真','Tradeoff'],['tabletop_relations','质量均好，物理更合理','PASS']],[48*mm,84*mm,38*mm])
    H('11. 官方 benchmark 使用边界');S.append(P('NVIDIA 当前 H20 PyTorch 720p/1 为 931.39s，但官方说明标准视频 workload 通常为189 frames；本项目为121 frames，且计时范围不同，因此不能直接计算931.39/344.612。正式加速比只使用本机同协议 B0/P4。'));S.append(P('<link href="https://github.com/NVIDIA/cosmos/blob/main/inference_benchmarks.md#text-to-video-t2v" color="#155EEF">NVIDIA Cosmos inference benchmark</link>'))
    H('12. 可复现命令与实现');cmd='env WARMUP=1 STEPS=35 RESIDUAL_CACHE_STEPS=11,13,15,17,19,21,23 RESIDUAL_CACHE_REL_L1_THRESHOLD=0.0545 RESIDUAL_CACHE_PROTECTED_FIRST_STEPS=5 RESIDUAL_CACHE_PROTECTED_LAST_STEPS=5 RESIDUAL_CACHE_PROTECTED_STEPS=25 bash /root/autodl-tmp/cosmos3_t2v_p4/run_residual_cache.sh';S.append(P(escape(cmd),ParagraphStyle('code',parent=small,backColor=colors.HexColor('#F2F4F7'),borderPadding=5)))
    for x in ['/root/autodl-tmp/cosmos-framework/cosmos_framework/inference/quantization.py','/root/autodl-tmp/cosmos-framework/cosmos_framework/model/vfm/utils/residual_cache.py','/root/autodl-tmp/cosmos-framework/cosmos_framework/model/vfm/omni_mot_model.py','/root/autodl-tmp/cosmos3_t2v_p4/run_residual_cache.sh']:S.append(P(escape(x),small))
    H('13. 服务器可验证产物（绝对路径）');TB([['Artifact','Absolute path']]+[[k,str(v)] for k,v in paths.items()],[38*mm,132*mm],6.1);H2('同噪声完整视频')
    for k,v in videos.items():S.append(P(f'<b>{escape(k)}</b><br/><link href="file://{escape(str(v))}" color="#155EEF">{escape(str(v))}</link>',small))
    H('14. 局限性与结论');bullets(['结果针对单张H20、batch1、121帧720p，其他硬件/shape需重新验证','FP8与动态cache跨进程不bit-exact，质量结论依赖人工审核','快速运动场景会自动更多刷新，因此P4收益随prompt变化','进一步大幅加速可能需要蒸馏或更激进近似，并承担新的质量风险']);S.append(P('<b>结论：</b>最终 P4 Balanced 在保持35-step UniPC和严格同噪声质量门的前提下实现1.3968×端到端加速，Anchor + 3/3高风险prompt全部通过。'))
    Doc(str(OUT/'COSMOS3_T2V_FINAL_REPORT.pdf')).build(S)

def canvas(title,kicker):
    im=Image.new('RGB',(1920,1080),C['navy']);d=ImageDraw.Draw(im);d.rectangle((0,0,1920,12),fill=C['green']);d.text((90,55),kicker,font=F(25),fill='#36D1DC');d.text((90,105),title,font=F(52),fill=C['white']);d.line((90,190,1830,190),fill=C['grid'],width=2);d.text((90,1030),'Cosmos3-Nano 16B · H20 · 720p/121f',font=F(21),fill=C['muted']);return im,d
def paste(im,src,box):
    pic=Image.open(src).convert('RGB');x1,y1,x2,y2=box;s=min((x2-x1)/pic.width,(y2-y1)/pic.height);pic=pic.resize((int(pic.width*s),int(pic.height*s)),Image.Resampling.LANCZOS);im.paste(pic,(x1+(x2-x1-pic.width)//2,y1+(y2-y1-pic.height)//2))
def bullets_slide(d,items,x=110,y=240,w=1660,size=34,gap=20):
    for t in items:d.ellipse((x,y+12,x+13,y+25),fill=C['green']);y=draw_lines(d,t,(x+35,y),F(size),C['white'],w-35,9)+gap

def build_ppt():
    S=[];links={}
    im,d=canvas('Cosmos3-Nano 16B 单卡 T2V 推理加速','FINAL REPORT');d.text((90,265),'720p · 121 frames · 24 FPS · 35-step UniPC',font=F(39),fill=C['muted'])
    for i,(v,l) in enumerate([('344.612s','最终E2E'),('1.3968×','vs B0'),('28.41%','时延下降'),('4/4','质量PASS')]):x=90+i*445;d.rounded_rectangle((x,420,x+395,650),radius=28,fill=C['navy2']);d.text((x+197,500),v,font=F(52),fill=C['green'],anchor='mm');d.text((x+197,590),l,font=F(27),fill='white',anchor='mm')
    d.text((90,760),'无训练 · 不减默认35步 · 同初始噪声人工质量门',font=F(39),fill='white');S.append(im)
    im,d=canvas('任务约束与实验协议','01 · OBJECTIVE');bullets_slide(d,['1× NVIDIA H20；1280×720；121帧；24 FPS。','35-step UniPC、shift10、CFG6；最终结果不依赖减步。','完全相同 initial noise；每个独立进程 warmup 后正式计时。','质量由完整同噪声视频人工审核；联系图不替代播放。','E2E、阶段 breakdown、GPU telemetry 和失败候选全部保留。'],y=245,size=36);S.append(im)
    im,d=canvas('Baseline：94% 时间在双 CFG MoT forward','02 · PROFILE');paste(im,GEN/'breakdown.png',(100,220,1820,910));d.text((100,945),'FA3 已启用；VAE decode 仅3.75%，优化重点明确。',font=F(30),fill='white');S.append(im)
    im,d=canvas('P1：没有收益的候选也要公开','03 · NEGATIVE RESULTS');bullets_slide(d,['Static compile：481.142s，1.00046×，噪声级。','CUDA Graph shape smoke：慢0.23%，VRAM +3.1GB。','Text K/V cache：2-step E2E仅1.0037×。','结论：大算子计算主导，launch/Python不是瓶颈。'],y=260,size=39,gap=30);S.append(im)
    im,d=canvas('P2：最后10步关闭 unconditional CFG','04 · DIFFUSION');bullets_slide(d,['Scheduler 35次、conditional 35次保持不变。','Unconditional 35→25；220.726→157.650s。','E2E 481.361→418.109s。','30/28-step只作消融，不进入最终主线。'],x=100,y=240,w=760,size=31);paste(im,GEN/'cfg.jpg',(900,225,1820,900));S.append(im)
    im,d=canvas('P3：只量化 generation 主干','05 · QUANTIZATION');bullets_slide(d,['第4–31层，140个q/o/MLP Linear。','动态activation + weight FP8 E4M3。','K/V、text、首尾层、Norm、RoPE、Attention、VAE保持BF16。','E2E median 359.067s；CV 0.0127%。'],x=100,y=240,w=790,size=30);paste(im,GEN/'fp8.jpg',(930,225,1820,900));S.append(im)
    im,d=canvas('P4：动态阈值 conditional residual cache','06 · FEATURE REUSE');paste(im,GEN/'cache.png',(80,220,1080,900));bullets_slide(d,['层8–27；候选11–23。','relative-L1≤0.0545才复用。','首尾5步和CFG切换点保护。','候选非相邻；失败完整刷新。','每轮命中5步，跳过100个block。'],x=1130,y=260,w=650,size=30);S.append(im)
    im,d=canvas('最终结果：1.3968×，时延下降28.41%','07 · PERFORMANCE');paste(im,GEN/'e2e.png',(80,215,1840,920));S.append(im)
    im,d=canvas('每项收益都能从 breakdown 归因','08 · BREAKDOWN');paste(im,GEN/'breakdown.png',(70,210,1850,920));S.append(im)
    im,d=canvas('三进程稳态复测：范围仅0.262秒','09 · STABILITY');paste(im,GEN/'stability.png',(90,225,1830,900));d.text((90,945),'边界step会变化，但命中数=5，E2E/VRAM稳定。',font=F(30),fill='white');S.append(im)
    im,d=canvas('Anchor 同噪声视觉对比','10 · QUALITY');paste(im,GEN/'p4.jpg',(90,220,1830,860));d.text((90,905),'左P3 / 右P4 Balanced · 完整视频人工审核 PASS',font=F(32),fill=C['green']);d.rounded_rectangle((90,960,560,1018),radius=14,fill=C['blue']);d.text((325,989),'点击打开完整MP4',font=F(25),fill='white',anchor='mm');S.append(im);links[len(S)-1]=str(videos['P3 vs P4 Balanced'])
    im,d=canvas('多 Prompt 质量门：Anchor + 3/3 PASS','11 · GENERALIZATION');paste(im,GEN/'quality.png',(80,220,1840,910));S.append(im)
    im,d=canvas('Aggressive 更快，但被质量门淘汰','12 · QUALITY BOUNDARY');paste(im,GEN/'aggr.jpg',(80,220,1160,900));bullets_slide(d,['Aggressive：335.881s。','仅比Balanced快约2.6%。','human_handoff边缘变模糊。','严格质量门FAIL，保留为边界消融。'],x=1210,y=300,w=600,size=31);S.append(im)
    im,d=canvas('官方 benchmark 只作背景参考','13 · REFERENCE');bullets_slide(d,['NVIDIA H20 PyTorch 720p/1：931.39s。','官方标准视频workload通常为189帧；本项目121帧。','计时范围也不同，不能直接计算2.70×。','正式加速比只用本机同协议B0 481.361s。'],y=280,size=40,gap=32);S.append(im)
    im,d=canvas('服务器可验证产物','14 · EVIDENCE');y=225
    for x in [str(paths['Exact initial noise']),str(paths['Baseline benchmark']),str(paths['P4 repeat benchmark']),str(paths['Quality manifest']),str(paths['P4 status'])]:d.rounded_rectangle((90,y,1830,y+125),radius=18,fill=C['navy2']);draw_lines(d,x,(120,y+25),F(25),'white',1660,6);y+=145
    S.append(im)
    im,d=canvas('结论','15 · TAKEAWAY');bullets_slide(d,['344.612s median；1.3968×；时延下降28.41%。','保持35-step UniPC；无训练、无蒸馏、无VAE修改。','有效三项：CFG裁剪、generation-selective FP8、dynamic residual reuse。','Anchor + 3/3高风险prompt同噪声人工质量PASS。','公开负结果与失败案例，结果可审计、可复现。'],y=245,size=38,gap=28);S.append(im)
    prs=Presentation();prs.slide_width=Inches(13.333333);prs.slide_height=Inches(7.5);blank=prs.slide_layouts[6];sd=GEN/'slides';sd.mkdir(exist_ok=True)
    for i,im in enumerate(S):
        fp=sd/f'slide_{i+1:02d}.png';im.save(fp,quality=96);sl=prs.slides.add_slide(blank);sl.shapes.add_picture(str(fp),0,0,width=prs.slide_width,height=prs.slide_height)
        if i in links:
            sh=sl.shapes.add_shape(1,Inches(.625),Inches(6.667),Inches(3.264),Inches(.403));sh.fill.background();sh.line.fill.background();sh.click_action.hyperlink.address='file://'+links[i]
    prs.save(OUT/'COSMOS3_T2V_PRESENTATION.pptx')

def main():
    assets();write_sources();build_pdf();build_ppt();print('done')
if __name__=='__main__':main()
